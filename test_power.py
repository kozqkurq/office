from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

# 演習
prs2 = Presentation("test.pptx")
title_slide_layout = prs2.slide_layouts[1]
slide = prs2.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "パワーポイント演習"
subtitle.text = "ピーナッツが好き"

prs2.save('test2.pptx')
count = 0
for slide in prs2.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                count += len(run.text)

print(count)